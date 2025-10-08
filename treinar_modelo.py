import os
import re
import pickle
<<<<<<< HEAD
import re
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.linear_model import LogisticRegression
from sklearn.model_selection import train_test_split
from sklearn.metrics import classification_report

# -----------------------
# Funções utilitárias
# -----------------------
def extrair_texto_pdf(path):
    texto = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            txt = page.extract_text()
            if not txt:
                # PDF como imagem, usa OCR
                imagem = page.to_image(resolution=300).original
                txt = pytesseract.image_to_string(imagem, lang="por")
            texto += txt + " "
    return texto

def extrair_texto_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def limpar_texto(texto):
    texto = texto.lower().strip()
    texto = re.sub(r'\s+', ' ', texto)
    texto = re.sub(r'[^\w\s]', '', texto)
    return texto

def extrair_texto_arquivo(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extrair_texto_pdf(path)
    elif ext == ".docx":
        return extrair_texto_docx(path)
    else:
        return None

# -----------------------
# Carregar dados
# -----------------------
# Estrutura: pasta "dados" com subpastas por vaga, contendo currículos
DATA_DIR = "dados"  # Substitua pelo caminho da sua pasta de dados
X, y = [], []

for vaga in os.listdir(DATA_DIR):
    pasta_vaga = os.path.join(DATA_DIR, vaga)
    if not os.path.isdir(pasta_vaga):
        continue
    for arquivo in os.listdir(pasta_vaga):
        path = os.path.join(pasta_vaga, arquivo)
        try:
            texto = extrair_texto_arquivo(path)
            if texto:
                texto = limpar_texto(texto)
                if texto.strip():
                    X.append(texto)
                    y.append(vaga.upper())  # Sempre maiúscula para consistência
        except Exception as e:
            print(f"Erro ao processar {path}: {e}")

# -----------------------
# Treinamento
# -----------------------
vectorizer = TfidfVectorizer()
X_vect = vectorizer.fit_transform(X)

clf = LogisticRegression(max_iter=500)
X_train, X_test, y_train, y_test = train_test_split(X_vect, y, test_size=0.2, random_state=42)
clf.fit(X_train, y_train)

# Avaliação
y_pred = clf.predict(X_test)
print(classification_report(y_test, y_pred))

# -----------------------
# Salvar modelo
# -----------------------
with open("modelo_curriculos.pkl", "wb") as f:
    pickle.dump((clf, vectorizer), f)

print("Treinamento concluído e modelo salvo em modelo_curriculos.pkl")
=======
import shutil
import logging
import unicodedata
from collections import Counter

import numpy as np
import pandas as pd

import pdfplumber
import docx
import pytesseract
from PIL import Image
from pptx import Presentation

import nltk
from nltk.corpus import stopwords

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.preprocessing import LabelEncoder, normalize
from sklearn.decomposition import PCA
from sklearn.model_selection import train_test_split, StratifiedKFold
from sklearn.feature_selection import SelectKBest, chi2
from sklearn.metrics import classification_report, f1_score

import xgboost as xgb
from imblearn.over_sampling import RandomOverSampler

from sentence_transformers import SentenceTransformer
from scipy.sparse import csr_matrix, hstack

import spacy

# --- Setup logging ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s: %(message)s')

# --- Preparar NLTK e spaCy ---
for res in ["stopwords", "punkt"]:
    try:
        nltk.data.find(f"corpora/{res}")
    except LookupError:
        nltk.download(res)

STOPWORDS = set(stopwords.words("portuguese"))

# Carregar o modelo spaCy para português (NER, POS etc)
# Você precisa instalar: `python -m spacy download pt_core_news_sm`
nlp_spacy = spacy.load("pt_core_news_sm")

def limpar_texto_avancado(texto: str) -> str:
    texto = unicodedata.normalize("NFKD", texto).encode("ASCII", "ignore").decode("utf-8")
    texto = texto.lower()
    texto = re.sub(r"\S+@\S+", " ", texto)
    texto = re.sub(r"\d+", " ", texto)
    texto = re.sub(r"[^a-z\s]", " ", texto)
    texto = re.sub(r"\s+", " ", texto).strip()
    palavras = [p for p in texto.split() if p not in STOPWORDS]
    return " ".join(palavras)

def extrair_texto_arquivo(filepath):
    """Extrai texto bruto (sem limpeza) de vários formatos."""
    ext = os.path.splitext(filepath)[1].lower()
    texto = ""
    try:
        if ext == ".pdf":
            with pdfplumber.open(filepath) as pdf:
                texto = " ".join([p.extract_text() or "" for p in pdf.pages])
            if not texto.strip():
                # fallback OCR por página
                with pdfplumber.open(filepath) as pdf:
                    for page in pdf.pages:
                        img = page.to_image(resolution=300).original
                        texto += pytesseract.image_to_string(img, lang="por")
        elif ext in (".docx", ".doc"):
            doc = docx.Document(filepath)
            texto = " ".join([p.text for p in doc.paragraphs])
        elif ext == ".txt":
            with open(filepath, encoding="utf-8", errors="ignore") as f:
                texto = f.read()
        elif ext == ".pptx":
            prs = Presentation(filepath)
            partes = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        partes.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            partes.append(" ".join(cell.text for cell in row.cells))
            texto = " ".join(partes)
        elif ext in (".xlsx", ".xls"):
            df = pd.read_excel(filepath, sheet_name=None)
            partes = []
            for aba in df.values():
                for _, row in aba.iterrows():
                    partes.append(" ".join(row.dropna().astype(str)))
            texto = " ".join(partes)
    except Exception as e:
        logging.warning(f"Falha ao extrair de {filepath}: {e}")
    return texto

# Heurísticas / regex extras
def extrair_formacoes(texto: str):
    """Busca expressões de formação acadêmica (ex: graduação em, bacharel, tecnólogo etc)."""
    padroes = [
        r"graduação em ([a-z\s]+)",
        r"bacharelado em ([a-z\s]+)",
        r"técnico em ([a-z\s]+)",
        r"licenciatura em ([a-z\s]+)",
        r"especialização em ([a-z\s]+)",
        r"mestrado em ([a-z\s]+)",
        r"doutorado em ([a-z\s]+)",
    ]
    found = []
    for p in padroes:
        matches = re.findall(p, texto)
        for m in matches:
            found.append(m.strip())
    return list(set(found))

def extrair_idiomas(texto: str):
    """Busca níveis de idiomas."""
    # termos comuns
    niveis = ["fluente", "avançado", "intermediário", "básico", "nível"]
    idiomas = ["inglês", "espanhol", "francês", "alemão", "italiano"]
    found = []
    for idi in idiomas:
        for niv in niveis:
            if f"{idi} {niv}" in texto:
                found.append(f"{idi} {niv}")
            elif idi in texto:
                found.append(idi)
    return list(set(found))

def extrair_certificacoes(texto: str):
    """Busca siglas comuns de certificações (ex: PMP, SCRUM, ITIL etc)."""
    siglas = ["PMP", "SCRUM", "ITIL", "Lean", "Six Sigma", "AWS", "Azure", "Google Cloud"]
    found = []
    for s in siglas:
        if s.lower() in texto.lower():
            found.append(s)
    return found

def extrair_anos_experiencia(texto: str):
    """Tentativa de extrair quantos anos de experiência ele menciona."""
    # padrões como 'X anos', 'anos de experiência'
    matches = re.findall(r"(\d+)\s+anos", texto)
    if matches:
        # pegar o maior número
        try:
            nums = [int(m) for m in matches]
            return max(nums)
        except:
            return None
    return None

def extrair_entidades_spacy(texto: str):
    """Usa spaCy NER para extrair entidades (ORG, PER, LOC)."""
    doc = nlp_spacy(texto)
    ent_map = {"ORG": [], "PERSON": [], "LOC": []}
    for ent in doc.ents:
        if ent.label_ in ent_map:
            ent_map[ent.label_].append(ent.text.strip())
    # deduplicar
    for k in ent_map:
        ent_map[k] = list(set(ent_map[k]))
    return ent_map

def extrair_features_estruturadas(texto: str):
    """Dividir o texto em seções heurísticas e extrair contagens / sinais."""
    seccoes = ["experiência", "formação", "educação", "certificações", "habilidades", "competências"]
    lower = texto.lower()
    # contar onde seções aparecem
    feat = {}
    for s in seccoes:
        feat[f"has_section_{s}"] = int(s in lower)
    # contagem de parágrafos / linhas
    feat["num_linhas"] = texto.count("\n") + 1
    feat["comprimento"] = len(texto)
    return feat

############################
# Pipeline de coleta de dados
############################

# Diretório com currículos organizados em pastas por área
pasta_base = r"C:\Users\Usuario\CURRICULOS"
texts = []
labels = []
metadatas = []  # para armazenar dados extras por candidato

for area in os.listdir(pasta_base):
    dir_area = os.path.join(pasta_base, area)
    if not os.path.isdir(dir_area):
        continue
    for fname in os.listdir(dir_area):
        fp = os.path.join(dir_area, fname)
        texto_bruto = extrair_texto_arquivo(fp)
        if not texto_bruto.strip():
            continue
        texto_limpo = limpar_texto_avancado(texto_bruto)
        texts.append(texto_limpo)
        labels.append(area.upper())
        # metadata
        meta = {
            "formacoes": extrair_formacoes(texto_bruto),
            "idiomas": extrair_idiomas(texto_bruto),
            "certificacoes": extrair_certificacoes(texto_bruto),
            "anos_exp": extrair_anos_experiencia(texto_bruto),
            "entidades": extrair_entidades_spacy(texto_bruto),
            **extrair_features_estruturadas(texto_bruto)
        }
        metadatas.append(meta)

logging.info(f"Total de textos: {len(texts)}")

# Vetorização TF-IDF + features de palavras-chave
word_vectorizer = TfidfVectorizer(analyzer="word", ngram_range=(1,2), max_features=5000)
char_vectorizer = TfidfVectorizer(analyzer="char", ngram_range=(2,5), max_features=5000)
Xw = word_vectorizer.fit_transform(texts)
Xc = char_vectorizer.fit_transform(texts)

# Carregar palavras-chave
planilha = pd.read_excel("palavras_chave.xlsx")
planilha.columns = planilha.columns.str.strip().str.upper()
palavras_chave_dict = planilha.groupby("AREA")["PALAVRA_CHAVE"].apply(list).to_dict()

def extrair_features_chave(texto: str):
    return [int(any(p.lower() in texto for p in palavras)) for palavras in palavras_chave_dict.values()]

X_chaves = csr_matrix([extrair_features_chave(t) for t in texts])

# Convert metadatas em features densas
df_meta = pd.DataFrame(metadatas)
# Exemplo: transformar formacoes / idiomas / certificações em contagens ou vetores binários
# Pode ser expandido conforme domínio
# Por exemplo, contar quantas formações = len(meta["formacoes"])
feat_formacoes = df_meta["formacoes"].apply(len).fillna(0).astype(int).values.reshape(-1, 1)
feat_idiomas = df_meta["idiomas"].apply(len).fillna(0).astype(int).values.reshape(-1, 1)
feat_cert = df_meta["certificacoes"].apply(len).fillna(0).astype(int).values.reshape(-1, 1)
feat_anos = df_meta["anos_exp"].fillna(0).astype(int).values.reshape(-1, 1)

# Outros metas: seção booleans + etc
feat_estruts = df_meta[[col for col in df_meta.columns if col.startswith("has_section_") or col in ("num_linhas","comprimento")]]
feat_estruts = feat_estruts.fillna(0).astype(int).values

# Concatena todas features densas
X_meta = np.hstack([feat_formacoes, feat_idiomas, feat_cert, feat_anos, feat_estruts])
# Normalizar meta features
from sklearn.preprocessing import StandardScaler
scaler_meta = StandardScaler()
X_meta_scaled = scaler_meta.fit_transform(X_meta)

# Embeddings semânticos
modelo_emb = SentenceTransformer('neuralmind/bert-base-portuguese-cased')
X_emb = modelo_emb.encode(texts, show_progress_bar=True)
X_emb = normalize(X_emb)

# Reduzir dimensões dos embeddings
from sklearn.decomposition import PCA
pca = PCA(n_components=50, random_state=42)
X_emb_reduz = pca.fit_transform(X_emb)

# Converter embeddings reduzidos para matriz esparsa para combinar com esparsas
X_emb_sparse = csr_matrix(X_emb_reduz)

# Combinar features esparsas (TF-IDF + chaves) + embeddings + metadata densa
X_tfidf_chaves = hstack([Xw, Xc, X_chaves])
selector = SelectKBest(chi2, k=min(5000, X_tfidf_chaves.shape[1]))
X_selected = selector.fit_transform(X_tfidf_chaves, labels)

# converter X_meta para esparso (csr) para empilhar
X_meta_sparse = csr_matrix(X_meta_scaled)

X_full = hstack([X_selected, X_emb_sparse, X_meta_sparse])
logging.info(f"Shape final de features: {X_full.shape}")

# Encode labels
le = LabelEncoder()
y = le.fit_transform(labels)
num_classes = len(le.classes_)

# Oversampling com todas as features
ros = RandomOverSampler(random_state=42)
X_bal, y_bal = ros.fit_resample(X_full, y)
logging.info(f"Apos oversampling: {X_bal.shape}, classes: {Counter(y_bal)}")

# Sample weights
counter = Counter(y_bal)
total = sum(counter.values())
class_weights = {i: total / (num_classes * cnt) for i, cnt in counter.items()}
sample_weights = np.array([class_weights[yval] for yval in y_bal])

# Separar treino/teste
X_train, X_test, y_train, y_test, sw_train, sw_test = train_test_split(
    X_bal, y_bal, sample_weights, test_size=0.2, random_state=42, stratify=y_bal
)

# Treinar XGBoost com early stopping (parâmetro no construtor)
clf = xgb.XGBClassifier(
    objective="multi:softprob",
    num_class=num_classes,
    eval_metric="mlogloss",
    n_estimators=700,
    max_depth=6,
    learning_rate=0.1,
    subsample=0.8,
    colsample_bytree=0.8,
    tree_method="hist",
    random_state=42,
    early_stopping_rounds=50
)

clf.fit(
    X_train, y_train,
    sample_weight=sw_train,
    eval_set=[(X_test, y_test)],
    verbose=10
)

# Avaliação no conjunto de teste
y_pred = clf.predict(X_test)
print("\n=== Relatório de classificação (conjunto de teste) ===\n")
print(classification_report(y_test, y_pred, target_names=le.classes_, zero_division=0))

# Validação cruzada (só métricas agregadas, sem imprimir relatório inteiro por fold)
cv = StratifiedKFold(n_splits=5, shuffle=True, random_state=42)
f1_scores = []
for tr_idx, te_idx in cv.split(X_bal, y_bal):
    Xt, Xv = X_bal[tr_idx], X_bal[te_idx]
    yt, yv = y_bal[tr_idx], y_bal[te_idx]
    swt = np.array([class_weights[y] for y in yt])

    model_cv = xgb.XGBClassifier(
        objective="multi:softprob",
        num_class=num_classes,
        eval_metric="mlogloss",
        n_estimators=700,
        max_depth=6,
        learning_rate=0.1,
        subsample=0.8,
        colsample_bytree=0.8,
        tree_method="hist",
        random_state=42,
        early_stopping_rounds=50
    )
    model_cv.fit(Xt, yt, sample_weight=swt, eval_set=[(Xv, yv)], verbose=False)
    yv_pred = model_cv.predict(Xv)
    f1_scores.append(f1_score(yv, yv_pred, average="weighted"))

media_f1 = np.mean(f1_scores)
logging.info(f"F1 médio na validação cruzada (5 folds): {media_f1:.4f}")

# Salvar modelo e artefatos
modelo_emb.save("modelo_emb_dir")
with open("modelo_curriculos_super_avancado.pkl", "wb") as f:
    pickle.dump({
        "clf": clf,
        "word_vectorizer": word_vectorizer,
        "char_vectorizer": char_vectorizer,
        "palavras_chave_dict": palavras_chave_dict,
        "selector": selector,
        "label_encoder": le,
        "pca": pca,
        "scaler_meta": scaler_meta,
        "emb_model_path": "modelo_emb_dir"
    }, f)

logging.info("Treinamento finalizado, modelo salvo.")
>>>>>>> 1f90bbe (Salvar alterações pendentes antes do rebase)
